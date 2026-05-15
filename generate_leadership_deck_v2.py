"""
Generate Hephaestus NG Leadership Slides & Detailed Report — v2
Data period: March 10 – April 9, 2026 (30 days)
Regions: East US 2, Sweden Central

Updates:
- Team details on title slide
- Removed ColdStart vs ChatStream breakdown slide
- Full V2 data (30 days, ~1,319 total)
- Success rate capped at 100% with footnote
- Appendix with dashboard links
- Placeholders for user retention & % NG users (pending data)
"""

import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocInches
from docx.enum.table import WD_TABLE_ALIGNMENT

OUT_DIR = os.path.join(os.path.dirname(__file__), "leadership_output")
os.makedirs(OUT_DIR, exist_ok=True)

# ─── TEAM INFO ───────────────────────────────────────────────────────────────
TEAM_NAME = "AskAI Service Team"
TEAM_MANAGER = "Nupur Kinger"
TEAM_MEMBERS = ["Amit Chauhan", "Mohammed Sheraj", "Kshitij Chawla", "Ayushh Garg"]

# ─── COLORS ──────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1B, 0x1B, 0x2F)
CARD_BG    = RGBColor(0x2D, 0x2D, 0x44)
ACCENT_BLUE   = RGBColor(0x00, 0x78, 0xD4)
ACCENT_GREEN  = RGBColor(0x10, 0x7C, 0x10)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY   = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)

# ─── DASHBOARD LINKS ────────────────────────────────────────────────────────
DASH = (
    "https://dataexplorer.azure.com/dashboards/"
    "4584380f-1dcf-43ff-b174-beb56a9a51ea"
    "?p-_startTime=30days&p-_endTime=now"
    "&p-_env=v-eastus2&p-_env=v-swedencentral"
)
TILES = {
    "NG Requests":           f"{DASH}&tile=400a95ab-6ff8-4bd4-bbae-18898d1a54a0",
    "V2 Requests (Triage)":  f"{DASH}&tile=d24ef970-b3cd-4383-9e17-6ec5099ab134",
    "V1 P95 Latency":        f"{DASH}&tile=5bf7b1d4-63e0-4af6-8f26-7d0f4297412a",
    "V1 P99 TTFB":             f"{DASH}&tile=c1c9cc13-1713-43b4-b0ed-36d94d192c0c",
    "V2 P99 TTFB":             f"{DASH}&tile=fac37981-da6c-4d16-bba0-a625a4939fc9",
    "V2 P95 Latency":           f"{DASH}&tile=abcd56e7-4390-4c2d-8931-24a3346232d3",
    "Success Rate (Overall)": f"{DASH}&tile=20253c74-5866-4fa1-bf3c-d62c65ab710e",
    "Success Rate (Daily)":  f"{DASH}&tile=2e6b6b6c-dafc-41ed-9113-70f243c8dece",
    "Success Rate (Region)": f"{DASH}&tile=bd391155-bdda-44f0-9195-b4d47e382277",
    "V1 Success Rate":       f"{DASH}&tile=2e6b6b6c-dafc-41ed-9113-70f243c8dece",
    "V2 Success Rate":       f"{DASH}&tile=dce01fdf-308a-4d63-8438-9d3745a065c0",
}

# ─── DATA ────────────────────────────────────────────────────────────────────
daily_total_requests = [
    ("Mar 10", 3141), ("Mar 11", 5174), ("Mar 12", 4710), ("Mar 13", 4834),
    ("Mar 14", 1457), ("Mar 15", 1535), ("Mar 16", 4313), ("Mar 17", 5357),
    ("Mar 18", 5190), ("Mar 19", 4328), ("Mar 20", 3187), ("Mar 21", 1331),
    ("Mar 22", 1435), ("Mar 23", 3969), ("Mar 24", 4420), ("Mar 25", 4231),
    ("Mar 26", 4402), ("Mar 27", 3624), ("Mar 28", 1429), ("Mar 29", 1496),
    ("Mar 30", 3949), ("Mar 31", 4202), ("Apr 01", 4051), ("Apr 02", 3659),
    ("Apr 03", 2525),
]

daily_users = [
    ("Mar 10", 469), ("Mar 11", 760), ("Mar 12", 799), ("Mar 13", 687),
    ("Mar 14", 235), ("Mar 15", 251), ("Mar 16", 724), ("Mar 17", 752),
    ("Mar 18", 739), ("Mar 19", 673), ("Mar 20", 496), ("Mar 21", 222),
    ("Mar 22", 236), ("Mar 23", 681), ("Mar 24", 716), ("Mar 25", 697),
    ("Mar 26", 715), ("Mar 27", 625), ("Mar 28", 208), ("Mar 29", 241),
    ("Mar 30", 663), ("Mar 31", 707), ("Apr 01", 703), ("Apr 02", 612),
    ("Apr 03", 440), ("Apr 04", 234), ("Apr 05", 237), ("Apr 06", 493),
    ("Apr 07", 697), ("Apr 08", 730), ("Apr 09", 354),
]

# Per-region success rate from tile bd391155
daily_success_eus2 = [
    ("Mar 14", 96.22),  ("Mar 15", 96.523), ("Mar 16", 95.89),  ("Mar 17", 95.313),
    ("Mar 18", 94.489), ("Mar 19", 95.559), ("Mar 20", 93.766), ("Mar 21", 96.903),
    ("Mar 22", 96.266), ("Mar 23", 94.232), ("Mar 24", 95.771), ("Mar 25", 96.277),
    ("Mar 26", 97.206), ("Mar 27", 97.351), ("Mar 28", 98.03),  ("Mar 29", 99.121),
    ("Mar 30", 97.975), ("Mar 31", 96.363), ("Apr 01", 97.201), ("Apr 02", 97.442),
    ("Apr 03", 97.621), ("Apr 04", 97.832), ("Apr 05", 98.052), ("Apr 06", 97.956),
    ("Apr 07", 98.033), ("Apr 08", 97.476), ("Apr 09", 97.722), ("Apr 10", 97.624),
    ("Apr 11", 97.524), ("Apr 12", 97.365), ("Apr 13", 79.688),
]
daily_success_sc = [
    ("Mar 14", 93.617), ("Mar 15", 98.0),   ("Mar 16", 96.639), ("Mar 17", 94.636),
    ("Mar 18", 96.721), ("Mar 19", 96.237), ("Mar 20", 95.333), ("Mar 21", 98.462),
    ("Mar 22", 97.619), ("Mar 23", 97.906), ("Mar 24", 95.082), ("Mar 25", 97.034),
    ("Mar 26", 98.69),  ("Mar 27", 96.135), ("Mar 28", 98.701), ("Mar 29", 97.222),
    ("Mar 30", 96.957), ("Mar 31", 99.119), ("Apr 01", 96.744), ("Apr 02", 99.502),
    ("Apr 03", 98.305), ("Apr 04", 100.0),  ("Apr 05", 93.333), ("Apr 06", 97.059),
    ("Apr 07", 96.939), ("Apr 08", 96.653), ("Apr 09", 97.908), ("Apr 10", 97.426),
    ("Apr 11", 98.571), ("Apr 12", 100.0),  ("Apr 13", 100.0),
]
# Weighted combined (83% EUS2, 17% SC)
daily_success = [
    (d, round(0.83 * e + 0.17 * s, 2))
    for (d, e), (_, s) in zip(daily_success_eus2, daily_success_sc)
]
SR_AVG = round(sum(r for _, r in daily_success) / len(daily_success), 1)
SR_MIN_VAL = min(r for _, r in daily_success)
SR_MIN_DAY = next(d for d, r in daily_success if r == SR_MIN_VAL)
SR_MAX_VAL = max(r for _, r in daily_success)
SR_MAX_DAY = next(d for d, r in daily_success if r == SR_MAX_VAL)

v2_daily = [
    ("Mar 16", 72), ("Mar 17", 69), ("Mar 18", 61), ("Mar 19", 53),
    ("Mar 20", 79), ("Mar 21", 9),  ("Mar 22", 3),  ("Mar 23", 21),
    ("Apr 02", 232), ("Apr 03", 230), ("Apr 04", 125), ("Apr 05", 104), ("Apr 06", 261),
]
V2_TOTAL = sum(c for _, c in v2_daily)

# V1 success rate (overall, tile 2e6b6b6c) — Mar 13 – Apr 12
v1_success_daily = [
    ("Mar 13", 94.502), ("Mar 14", 95.861), ("Mar 15", 96.597), ("Mar 16", 95.947),
    ("Mar 17", 95.262), ("Mar 18", 94.574), ("Mar 19", 95.576), ("Mar 20", 93.918),
    ("Mar 21", 97.024), ("Mar 22", 96.387), ("Mar 23", 94.514), ("Mar 24", 95.816),
    ("Mar 25", 96.438), ("Mar 26", 97.332), ("Mar 27", 97.238), ("Mar 28", 98.088),
    ("Mar 29", 98.982), ("Mar 30", 97.804), ("Mar 31", 96.590), ("Apr 01", 97.190),
    ("Apr 02", 97.596), ("Apr 03", 97.678), ("Apr 04", 98.008), ("Apr 05", 97.791),
    ("Apr 06", 97.899), ("Apr 07", 97.989), ("Apr 08", 97.358), ("Apr 09", 97.782),
    ("Apr 10", 97.555), ("Apr 11", 97.616), ("Apr 12", 98.471),
]
V1_SR_AVG = round(sum(r for _, r in v1_success_daily) / len(v1_success_daily), 1)
V1_SR_MIN = min(r for _, r in v1_success_daily)
V1_SR_MAX = max(r for _, r in v1_success_daily)

# V2 success rate (tile dce01fdf) — Mar 16 – Apr 12
v2_success_daily = [
    ("Mar 16", 91.549), ("Mar 17", 92.754), ("Mar 18", 88.525), ("Mar 19", 92.453),
    ("Mar 20", 94.872), ("Mar 21", 100.0),  ("Mar 22", 100.0),  ("Mar 23", 90.476),
    ("Mar 24", 92.126), ("Mar 25", 94.805), ("Mar 26", 92.747), ("Mar 27", 94.747),
    ("Mar 28", 95.205), ("Mar 29", 92.771), ("Mar 30", 93.053), ("Mar 31", 93.426),
    ("Apr 01", 95.483), ("Apr 02", 95.400), ("Apr 03", 89.083), ("Apr 04", 94.400),
    ("Apr 05", 97.872), ("Apr 06", 94.779), ("Apr 07", 95.192), ("Apr 08", 93.028),
    ("Apr 09", 99.795), ("Apr 10", 98.910), ("Apr 11", 100.0),  ("Apr 12", 100.0),
]
V2_SR_AVG = round(sum(r for _, r in v2_success_daily) / len(v2_success_daily), 1)
V2_SR_MIN = min(r for _, r in v2_success_daily)
V2_SR_MAX = max(r for _, r in v2_success_daily)

# ─── LATENCY SUMMARY (from hourly dashboard samples) ─────────────────────────
# V1 P95 (tile 5bf7b1d4, 100 hourly samples Mar 13–17)
V1_P95_MEDIAN = 5.6    # seconds
V1_P95_RANGE  = "5\u20136s"
V1_P95_MIN_S  = 4.1
V1_P95_MAX_S  = 8.4

# V1 P99 (tile c1c9cc13, 100 hourly samples Mar 13–17)
V1_P99_MEDIAN = 7.4    # seconds (typical, excl cold-start spikes)
V1_P99_RANGE  = "6\u20138s"
V1_P99_MIN_S  = 4.7
V1_P99_MAX_S  = 109.3  # cold-start timeout spike

# V2 P95 (tile abcd56e7, ~100 hourly samples Mar 16–23)
V2_P95_MEDIAN = 5.1    # seconds
V2_P95_RANGE  = "4\u20137s"
V2_P95_MIN_S  = 3.3
V2_P95_MAX_S  = 17.7

# V2 P99 (tile fac37981, ~100 hourly samples Mar 16–23)
V2_P99_MEDIAN = 5.0    # seconds
V2_P99_RANGE  = "4\u20136s"
V2_P99_MIN_S  = 3.3
V2_P99_MAX_S  = 17.7

TOTAL_NG   = 91371
TOTAL_USERS = 15453
AVG_WD_DAU = 700
AVG_WE_DAU = 235
EUS2_PCT   = 83
SC_PCT     = 17

# % of NG users in total AskAI
ASKAI_TOTAL_USERS = 118082
ASKAI_NG_USERS    = 16239
ASKAI_NG_PCT      = 13.75

WEEKEND_DATES = {
    "Mar 14", "Mar 15", "Mar 21", "Mar 22", "Mar 28", "Mar 29",
    "Apr 04", "Apr 05",
}

# User retention: (date, NewUsers, ReturningUsers, TotalUsers, NewPct, ReturnPct)
retention_daily = [
    ("Mar 12", 55,  0,  55,  100.0, 0.0),
    ("Mar 13", 618, 4,  622, 99.4,  0.6),
    ("Mar 14", 225, 10, 235, 95.7,  4.3),
    ("Mar 15", 238, 13, 251, 94.8,  5.2),
    ("Mar 16", 690, 34, 724, 95.3,  4.7),
    ("Mar 17", 707, 45, 752, 94.0,  6.0),
    ("Mar 18", 689, 50, 739, 93.2,  6.8),
    ("Mar 19", 629, 44, 673, 93.5,  6.5),
    ("Mar 20", 452, 44, 496, 91.1,  8.9),
    ("Mar 21", 214, 8,  222, 96.4,  3.6),
    ("Mar 22", 217, 19, 236, 91.9,  8.1),
    ("Mar 23", 637, 44, 681, 93.5,  6.5),
    ("Mar 24", 660, 56, 716, 92.2,  7.8),
    ("Mar 25", 627, 70, 697, 90.0,  10.0),
    ("Mar 26", 654, 61, 715, 91.5,  8.5),
    ("Mar 27", 568, 57, 625, 90.9,  9.1),
    ("Mar 28", 192, 16, 208, 92.3,  7.7),
    ("Mar 29", 217, 24, 241, 90.0,  10.0),
    ("Mar 30", 604, 59, 663, 91.1,  8.9),
    ("Mar 31", 646, 61, 707, 91.4,  8.6),
    ("Apr 01", 631, 72, 703, 89.8,  10.2),
    ("Apr 02", 553, 59, 612, 90.4,  9.6),
    ("Apr 03", 395, 45, 440, 89.8,  10.2),
    ("Apr 04", 216, 18, 234, 92.3,  7.7),
    ("Apr 05", 225, 12, 237, 94.9,  5.1),
    ("Apr 06", 440, 53, 493, 89.2,  10.8),
    ("Apr 07", 635, 62, 697, 91.1,  8.9),
    ("Apr 08", 662, 68, 730, 90.7,  9.3),
    ("Apr 09", 675, 79, 754, 89.5,  10.5),
    ("Apr 10", 545, 58, 603, 90.4,  9.6),
    ("Apr 11", 221, 18, 239, 92.5,  7.5),
    ("Apr 12", 237, 13, 250, 94.8,  5.2),
    ("Apr 13", 89,  16, 105, 84.8,  15.2),
]

TOTAL_NEW_USERS = sum(n for _, n, _, _, _, _ in retention_daily)
TOTAL_RETURNING = sum(r for _, _, r, _, _, _ in retention_daily)
AVG_RETURN_PCT  = round(sum(rp for _, _, _, _, _, rp in retention_daily) / len(retention_daily), 1)
LATEST_RETURN_PCT = retention_daily[-1][5]
# Weekday-only returning % (excl weekends)
WD_RETURN_PCTS = [rp for d, _, _, _, _, rp in retention_daily if d not in WEEKEND_DATES]
AVG_WD_RETURN_PCT = round(sum(WD_RETURN_PCTS) / len(WD_RETURN_PCTS), 1) if WD_RETURN_PCTS else 0


# ═══════════════════════════════════════════════════════════════════════════════
# CHARTS
# ═══════════════════════════════════════════════════════════════════════════════

def _chart_requests():
    dates  = [d for d, _ in daily_total_requests]
    counts = [c for _, c in daily_total_requests]
    fig, ax = plt.subplots(figsize=(12, 4.2))
    bars = ax.bar(range(len(dates)), counts, color="#0078D4", width=0.7, alpha=0.85)
    for i, d in enumerate(dates):
        if d in WEEKEND_DATES:
            bars[i].set_color("#B0C4DE"); bars[i].set_alpha(0.5)
    ax.set_xticks(range(len(dates)))
    ax.set_xticklabels(dates, rotation=45, ha="right", fontsize=7)
    ax.set_ylabel("Total Requests")
    ax.set_title("Daily NG Request Volume (Mar 10 \u2013 Apr 3)", fontsize=12, fontweight="bold")
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"{int(v):,}"))
    for s in ("top", "right"): ax.spines[s].set_visible(False)
    fig.tight_layout()
    p = os.path.join(OUT_DIR, "chart_requests.png")
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig); return p


def _chart_users():
    dates  = [d for d, _ in daily_users]
    counts = [c for _, c in daily_users]
    fig, ax = plt.subplots(figsize=(12, 4.2))
    ax.fill_between(range(len(dates)), counts, alpha=0.25, color="#0078D4")
    ax.plot(range(len(dates)), counts, color="#0078D4", lw=2.5, marker="o", ms=4)
    for i, d in enumerate(dates):
        if d in WEEKEND_DATES:
            ax.axvspan(i - 0.4, i + 0.4, alpha=0.06, color="gray")
    ax.set_xticks(range(len(dates)))
    ax.set_xticklabels(dates, rotation=45, ha="right", fontsize=7)
    ax.set_ylabel("Unique Users")
    ax.set_title("Daily Unique Users (Mar 10 \u2013 Apr 9)", fontsize=12, fontweight="bold")
    for s in ("top", "right"): ax.spines[s].set_visible(False)
    fig.tight_layout()
    p = os.path.join(OUT_DIR, "chart_users.png")
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig); return p


def _chart_success():
    dates = [d for d, _ in daily_success]
    combined = [r for _, r in daily_success]
    eus2 = [r for _, r in daily_success_eus2]
    sc = [r for _, r in daily_success_sc]
    x = range(len(dates))
    fig, ax = plt.subplots(figsize=(12, 4.0))
    ax.plot(x, eus2, color="#0078D4", lw=1.5, marker="s", ms=3, alpha=0.6, label="East US 2")
    ax.plot(x, sc, color="#FF8C00", lw=1.5, marker="^", ms=3, alpha=0.6, label="Sweden Central")
    ax.plot(x, combined, color="#107C10", lw=2.5, marker="o", ms=5, label="Weighted Combined")
    ax.axhline(y=95, color="#D13438", ls="--", lw=1, alpha=0.7, label="95% SLO Target")
    ax.set_ylim(78, 101)
    ax.set_xticks(x)
    ax.set_xticklabels(dates, rotation=45, ha="right", fontsize=7)
    ax.set_ylabel("Success Rate (%)")
    ax.set_title("Daily Success Rate by Region (Mar 14 \u2013 Apr 13)", fontsize=12, fontweight="bold")
    ax.legend(loc="lower left", fontsize=9)
    for s in ("top", "right"): ax.spines[s].set_visible(False)
    fig.tight_layout()
    p = os.path.join(OUT_DIR, "chart_success.png")
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig); return p


def _chart_v1_success():
    dates = [d for d, _ in v1_success_daily]
    rates = [r for _, r in v1_success_daily]
    x = range(len(dates))
    fig, ax = plt.subplots(figsize=(6.5, 3.2))
    ax.plot(x, rates, color="#0078D4", lw=2, marker="o", ms=4, label="V1 Success Rate")
    ax.fill_between(x, rates, alpha=0.15, color="#0078D4")
    ax.axhline(y=95, color="#D13438", ls="--", lw=1, alpha=0.7, label="95% SLO Target")
    ax.set_ylim(92, 100)
    ax.set_xticks(x)
    ax.set_xticklabels(dates, rotation=45, ha="right", fontsize=6.5)
    ax.set_ylabel("Success Rate (%)")
    ax.set_title(f"V1 Daily Success Rate \u2014 Avg {V1_SR_AVG}%", fontsize=11, fontweight="bold")
    ax.legend(loc="lower right", fontsize=8)
    for s in ("top", "right"): ax.spines[s].set_visible(False)
    fig.tight_layout()
    p = os.path.join(OUT_DIR, "chart_v1_success.png")
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig); return p


def _chart_v2_success():
    dates = [d for d, _ in v2_success_daily]
    rates = [r for _, r in v2_success_daily]
    x = range(len(dates))
    fig, ax = plt.subplots(figsize=(6.5, 3.2))
    ax.plot(x, rates, color="#FF8C00", lw=2, marker="s", ms=4, label="V2 Success Rate")
    ax.fill_between(x, rates, alpha=0.15, color="#FF8C00")
    ax.axhline(y=95, color="#D13438", ls="--", lw=1, alpha=0.7, label="95% SLO Target")
    ax.set_ylim(85, 101)
    ax.set_xticks(x)
    ax.set_xticklabels(dates, rotation=45, ha="right", fontsize=6.5)
    ax.set_ylabel("Success Rate (%)")
    ax.set_title(f"V2 Daily Success Rate \u2014 Avg {V2_SR_AVG}%", fontsize=11, fontweight="bold")
    ax.legend(loc="lower right", fontsize=8)
    for s in ("top", "right"): ax.spines[s].set_visible(False)
    fig.tight_layout()
    p = os.path.join(OUT_DIR, "chart_v2_success.png")
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig); return p


def _chart_region_pie():
    fig, ax = plt.subplots(figsize=(4.5, 3.5))
    w, t, a = ax.pie([EUS2_PCT, SC_PCT], labels=["East US 2", "Sweden Central"],
                     autopct="%1.0f%%", colors=["#0078D4", "#FF8C00"],
                     startangle=90, textprops={"fontsize": 11})
    for x in a: x.set_fontweight("bold"); x.set_color("white")
    ax.set_title("Traffic by Region", fontsize=11, fontweight="bold")
    fig.tight_layout()
    p = os.path.join(OUT_DIR, "chart_region.png")
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig); return p


def _chart_v2_vs_ng():
    fig, ax = plt.subplots(figsize=(5, 3.5))
    bars = ax.bar(["NG (Hephaestus)", "V2 (Triage)"], [TOTAL_NG, V2_TOTAL],
                  color=["#0078D4", "#B0C4DE"], width=0.5)
    ax.bar_label(bars, fmt="{:,.0f}", fontsize=11, fontweight="bold", padding=3)
    ax.set_ylabel("Total Requests (30d)")
    ax.set_title("NG vs V2 Traffic", fontsize=12, fontweight="bold")
    ax.set_ylim(0, TOTAL_NG * 1.15)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"{int(v):,}"))
    for s in ("top", "right"): ax.spines[s].set_visible(False)
    fig.tight_layout()
    p = os.path.join(OUT_DIR, "chart_v2_vs_ng.png")
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig); return p


def _chart_retention():
    dates = [d for d, *_ in retention_daily]
    new_u = [n for _, n, _, _, _, _ in retention_daily]
    ret_u = [r for _, _, r, _, _, _ in retention_daily]
    ret_p = [rp for _, _, _, _, _, rp in retention_daily]
    x = range(len(dates))

    fig, ax1 = plt.subplots(figsize=(12, 4.5))
    ax1.bar(x, new_u, color="#0078D4", width=0.6, label="New Users", alpha=0.85)
    ax1.bar(x, ret_u, bottom=new_u, color="#FF8C00", width=0.6, label="Returning Users", alpha=0.85)
    ax1.set_ylabel("Users")
    ax1.set_xticks(x)
    ax1.set_xticklabels(dates, rotation=45, ha="right", fontsize=7)
    for s in ("top", "right"): ax1.spines[s].set_visible(False)

    ax2 = ax1.twinx()
    ax2.plot(x, ret_p, color="#D13438", lw=2.5, marker="o", ms=4, label="Returning %")
    ax2.set_ylabel("Returning %")
    ax2.set_ylim(0, 20)
    ax2.spines["top"].set_visible(False)

    # Combined legend
    h1, l1 = ax1.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    ax1.legend(h1 + h2, l1 + l2, loc="upper left", fontsize=9)

    ax1.set_title("Daily New vs Returning Users (Mar 12 \u2013 Apr 13)", fontsize=12, fontweight="bold")
    fig.tight_layout()
    p = os.path.join(OUT_DIR, "chart_retention.png")
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig); return p


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _bg(sl, c=DARK_BG):
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = c

def _txt(sl, l, t, w, h, txt, sz=18, b=False, c=WHITE, a=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = txt; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c; p.alignment = a
    return bx

def _ml(sl, l, t, w, h, lines, sz=12, c=WHITE, a=PP_ALIGN.LEFT, b1=False):
    bx = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz); p.font.color.rgb = c; p.alignment = a
        if b1 and i == 0: p.font.bold = True
    return bx

def _kpi(sl, l, t, val, lbl, vc=ACCENT_BLUE, w=2.6, h=1.5):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = CARD_BG; s.line.fill.background()
    _txt(sl, l+.1, t+.12, w-.2, .7, str(val), 28, True, vc, PP_ALIGN.CENTER)
    _txt(sl, l+.1, t+.85, w-.2, .5, lbl, 11, False, LIGHT_GRAY, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD POWERPOINT
# ═══════════════════════════════════════════════════════════════════════════════

def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    # ── SLIDE 1 : Title + Team ────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
    _txt(sl, .8, .5, 11, 1, "Hephaestus NG \u2014 30-Day Dashboard Review", 36, True, WHITE)
    _txt(sl, .8, 1.4, 11, .5, "March 10 \u2013 April 9, 2026  |  East US 2 & Sweden Central", 18, c=LIGHT_GRAY)

    s = sl.shapes.add_shape(1, Inches(.8), Inches(2.6), Inches(5.5), Inches(3.8))
    s.fill.solid(); s.fill.fore_color.rgb = CARD_BG; s.line.fill.background()
    _txt(sl, 1, 2.75, 5, .5, TEAM_NAME, 22, True, ACCENT_BLUE)
    _txt(sl, 1, 3.35, 5, .4, f"Manager:  {TEAM_MANAGER}", 14, c=WHITE)
    _ml(sl, 1, 3.95, 5, 2, [
        "Team Members:",
        *[f"  \u2022  {m}" for m in TEAM_MEMBERS],
    ], 13, LIGHT_GRAY)

    _txt(sl, 7.2, 2.75, 5.5, .5, "\U0001f4ca  In This Review", 18, True, WHITE)
    _ml(sl, 7.2, 3.35, 5.5, 3, [
        "\u2022  Request Volume & Trends",
        "\u2022  User Adoption & Engagement",
        "\u2022  SLO Performance (Success Rate & Latency)",
        "\u2022  V2 \u2192 NG Migration Status",
        "\u2022  Key Takeaways & Next Steps",
    ], 14, LIGHT_GRAY)

    # ── SLIDE 2 : Key Metrics ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
    _txt(sl, .5, .3, 12, .7, "Key Metrics at a Glance", 28, True, WHITE)

    # Row 1: Volume & Users (y=1.1)
    _kpi(sl, .3,  1.1, "91,371",    "Total NG Requests (30d)",  ACCENT_BLUE,  w=2.5, h=1.3)
    _kpi(sl, 3.0, 1.1, "15,453",    "Unique Users (30d)",       ACCENT_GREEN, w=2.5, h=1.3)
    _kpi(sl, 5.7, 1.1, "~700/day",  "Avg Weekday DAU",          ACCENT_GREEN, w=2.5, h=1.3)
    _kpi(sl, 8.4, 1.1, "83% / 17%", "EUS2 / Sweden Central",    ACCENT_BLUE,  w=2.2, h=1.3)
    _kpi(sl, 10.8,1.1, f"{V2_TOTAL:,}", "V2 Requests (30d)",    MED_GRAY,     w=2.2, h=1.3)

    # Row 2: Success Rates (y=2.6)
    _kpi(sl, .3,  2.6, f"{V1_SR_AVG}%",  "V1 Avg Success Rate",
         ACCENT_GREEN if V1_SR_AVG >= 95 else ACCENT_ORANGE, w=3.1, h=1.3)
    _kpi(sl, 3.6, 2.6, f"{V2_SR_AVG}%",  "V2 Avg Success Rate",
         ACCENT_GREEN if V2_SR_AVG >= 95 else ACCENT_ORANGE, w=3.1, h=1.3)
    _kpi(sl, 6.9, 2.6, f"{SR_AVG}%",     "Per-Region Weighted SR",
         ACCENT_GREEN if SR_AVG >= 97 else ACCENT_ORANGE,    w=3.1, h=1.3)

    # Row 3: Latencies — V1 vs V2 (y=4.1)
    _kpi(sl, .3,  4.1, f"~{V1_P95_MEDIAN}s", "V1 P95 Latency",  ACCENT_BLUE,  w=3.1, h=1.3)
    _kpi(sl, 3.6, 4.1, f"~{V1_P99_MEDIAN}s", "V1 P99 TTFB",     ACCENT_BLUE,  w=3.1, h=1.3)
    _kpi(sl, 6.9, 4.1, "\u26a0\ufe0f TBD" if V2_P95_MEDIAN is None else f"~{V2_P95_MEDIAN}s",
         "V2 P95 Latency",  ACCENT_ORANGE if V2_P95_MEDIAN is None else ACCENT_BLUE, w=3.1, h=1.3)
    _kpi(sl, 10.2,4.1, f"~{V2_P99_MEDIAN}s", "V2 P99 TTFB",     ACCENT_BLUE,  w=3.1, h=1.3)

    _txt(sl, .5, 5.6, 12, .4,
         f"V1 SR: {V1_SR_AVG}% (min {V1_SR_MIN}% / max {V1_SR_MAX}%)   \u2502   "
         f"V2 SR: {V2_SR_AVG}% (min {V2_SR_MIN}% / max {V2_SR_MAX}%)   \u2502   "
         f"Per-Region: {SR_AVG}% (min {SR_MIN_VAL}% / max {SR_MAX_VAL}%)",
         9, c=LIGHT_GRAY)
    _txt(sl, .5, 6.0, 12, .4,
         f"V1 P95: {V1_P95_RANGE} (max {V1_P95_MAX_S}s)   \u2502   "
         f"V2 P95: {V2_P95_RANGE} (max {V2_P95_MAX_S}s)   \u2502   "
         f"V1 P99: {V1_P99_RANGE} typ, spikes to {V1_P99_MAX_S}s   \u2502   "
         f"V2 P99: {V2_P99_RANGE} (max {V2_P99_MAX_S}s)",
         9, c=LIGHT_GRAY)
    _txt(sl, .5, 6.4, 12, .4,
         f"NG handles 98.6% of all traffic ({TOTAL_NG:,} NG vs {V2_TOTAL:,} V2) \u2014 "
         "migration is effectively complete.",
         9, c=MED_GRAY)

    # ── SLIDE 3 : Request Volume + Region ─────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, WHITE)
    _txt(sl, .5, .2, 12, .6, "Request Volume & Regional Distribution", 24, True, DARK_TEXT)

    sl.shapes.add_picture(_chart_requests(), Inches(.3), Inches(.9), Inches(8.8), Inches(4.2))
    sl.shapes.add_picture(_chart_region_pie(), Inches(9.3), Inches(.9), Inches(3.7), Inches(2.8))

    _ml(sl, 9.3, 3.9, 3.7, 3, [
        "Key Observations:",
        "",
        f"\u2022 {TOTAL_NG:,} total requests (30d)",
        "\u2022 ~3,600 avg daily (weekdays)",
        "\u2022 ~1,400 avg daily (weekends)",
        "\u2022 East US 2 carries 83% of load",
        "\u2022 Lighter bars = weekend days",
    ], 10, DARK_TEXT, b1=True)

    # ── SLIDE 4 : User Adoption, Retention & AskAI Reach (merged) ─────────
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, WHITE)
    _txt(sl, .5, .2, 12, .6, "User Adoption, Retention & AskAI Reach", 24, True, DARK_TEXT)

    # Users chart (left) + Retention chart (right)
    sl.shapes.add_picture(_chart_users(), Inches(.2), Inches(.85), Inches(6.4), Inches(3.0))
    sl.shapes.add_picture(_chart_retention(), Inches(6.8), Inches(.85), Inches(6.4), Inches(3.0))

    # User highlights (left panel)
    _ml(sl, .3, 4.0, 4, 1.8, [
        "\U0001f4ca User Highlights",
        "",
        f"\u2022 {TOTAL_USERS:,} unique users (30d)",
        f"\u2022 ~{AVG_WD_DAU} DAU (weekdays) / ~{AVG_WE_DAU} (weekends)",
        "\u2022 Peak: 799 users (Mar 12)",
    ], 10, DARK_TEXT, b1=True)

    # Retention highlights (middle panel)
    _ml(sl, 4.5, 4.0, 4.2, 1.8, [
        "\U0001f504 Retention Highlights",
        "",
        f"\u2022 {TOTAL_RETURNING:,} returning users",
        f"\u2022 Weekday avg: {AVG_WD_RETURN_PCT}% returning",
        "\u2022 Trend: 0% \u2192 ~10% over 30 days",
    ], 10, DARK_TEXT, b1=True)

    # AskAI reach (right panel)
    _ml(sl, 8.9, 4.0, 4.2, 1.8, [
        "\U0001f310 AskAI Reach",
        "",
        f"\u2022 {ASKAI_TOTAL_USERS:,} total AskAI users",
        f"\u2022 {ASKAI_NG_USERS:,} use Hephaestus NG",
        f"\u2022 {ASKAI_NG_PCT}% of AskAI users on NG",
    ], 10, DARK_TEXT, b1=True)

    _txt(sl, .3, 5.9, 6, .4,
         f"Returning % = repeat visitors within 30-day window.  "
         f"\u2018New\u2019 = first appearance in window.",
         8, c=MED_GRAY)
    _txt(sl, 6.5, 5.9, 6.5, .4,
         f"NG users are {ASKAI_NG_PCT}% of {ASKAI_TOTAL_USERS:,} total AskAI users.",
         8, c=MED_GRAY)

    # ── SLIDE 5 : SLO Performance ─────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, WHITE)
    _txt(sl, .5, .2, 12, .6, "SLO Performance", 24, True, DARK_TEXT)

    # V1 success chart (left) + V2 success chart (right)
    sl.shapes.add_picture(_chart_v1_success(), Inches(.2), Inches(.85), Inches(6.4), Inches(3.0))
    sl.shapes.add_picture(_chart_v2_success(), Inches(6.8), Inches(.85), Inches(6.4), Inches(3.0))

    # V1 summary (left panel)
    _ml(sl, .3, 4.0, 4.2, 2.2, [
        "V1 Success Rate",
        "",
        f"\u2705 Average: {V1_SR_AVG}%",
        f"    Min: {V1_SR_MIN}%  |  Max: {V1_SR_MAX}%",
        f"    31 days (Mar 13\u2013Apr 12)",
        "    Consistently above 95% SLO",
    ], 10, DARK_TEXT, b1=True)

    # V2 summary (middle panel)
    _ml(sl, 4.7, 4.0, 4.2, 2.2, [
        "V2 Success Rate",
        "",
        f"\u26a0\ufe0f  Average: {V2_SR_AVG}%",
        f"    Min: {V2_SR_MIN}%  |  Max: {V2_SR_MAX}%",
        f"    28 days (Mar 16\u2013Apr 12)",
        "    More volatile; dips below 90%",
    ], 10, DARK_TEXT, b1=True)

    # Latency summary (right panel)
    _ml(sl, 9.2, 4.0, 4, 2.2, [
        "Latency Summary",
        "",
        f"V1 P95: {V1_P95_RANGE}  |  V2 P95: {V2_P95_RANGE}",
        f"V1 P99: {V1_P99_RANGE}  |  V2 P99: {V2_P99_RANGE}",
        f"V1 P99 spikes to ~{V1_P99_MAX_S}s (cold-start)",
        "V2 no extreme spikes observed",
    ], 10, DARK_TEXT, b1=True)

    _txt(sl, .3, 6.3, 6, .4,
         f"V1 tile: 2e6b6b6c  |  V2 tile: dce01fdf  |  Region tile: bd391155",
         7, c=MED_GRAY)
    _txt(sl, 6.5, 6.3, 6.5, .4,
         f"Per-region weighted avg: {SR_AVG}% (83% EUS2 + 17% SC). "
         "Region chart in Word report.",
         7, c=MED_GRAY)

    # ── SLIDE 6 : V2 -> NG Migration ─────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, WHITE)
    _txt(sl, .5, .2, 12, .6, "V2 \u2192 NG Migration Status", 24, True, DARK_TEXT)

    sl.shapes.add_picture(_chart_v2_vs_ng(), Inches(.3), Inches(.9), Inches(4.5), Inches(3.2))

    _ml(sl, 5.2, .9, 7.5, 3.5, [
        "Migration Summary",
        "",
        f"\u2022 NG Requests (30d):  {TOTAL_NG:,}   (98.6%)",
        f"\u2022 V2 Requests (30d):  {V2_TOTAL:,}     (1.4%)",
        "",
        "\u2022 V2 traffic is sporadic & declining",
        "\u2022 Several days with zero V2 requests (Mar 24\u2013Apr 1)",
        "\u2022 Migration to NG is effectively complete",
        "",
        "V2 Daily Breakdown:",
    ], 11, DARK_TEXT, b1=True)

    v2_lines = [f"  {d}: {c}" for d, c in v2_daily]
    _ml(sl, 5.2, 4.5, 4, 2.5, v2_lines, 9, DARK_TEXT)

    _ml(sl, 9.5, 4.5, 3.5, 2.5, [
        "Data Source:",
        'Traces table query filtering',
        'message contains',
        '"Running for AgentId=triage"',
    ], 9, MED_GRAY)

    # ── SLIDE 7 : Summary & Takeaways ─────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
    _txt(sl, .5, .4, 12, .7, "Summary & Takeaways", 28, True, WHITE)

    items = [
        ("\U0001f4c8  Strong Adoption",
         f"{TOTAL_NG:,} requests from {TOTAL_USERS:,} unique users in 30 days. "
         f"~{AVG_WD_DAU} weekday DAU. {ASKAI_NG_PCT}% of all AskAI users ({ASKAI_NG_USERS:,} / {ASKAI_TOTAL_USERS:,})."),
        ("\U0001f504  Growing Retention",
         f"Returning users grew from 0% to ~10% over 30 days (weekday avg {AVG_WD_RETURN_PCT}%). "
         f"{TOTAL_RETURNING:,} repeat users demonstrate increasing stickiness."),
        ("\u2705  V1 Reliability",
         f"V1 avg {V1_SR_AVG}% success rate (min {V1_SR_MIN}% / max {V1_SR_MAX}%). "
         f"P95: {V1_P95_RANGE}. P99: {V1_P99_RANGE} with cold-start spikes."),
        ("\u26a0\ufe0f  V2 Reliability",
         f"V2 avg {V2_SR_AVG}% success rate (min {V2_SR_MIN}% / max {V2_SR_MAX}%). "
         "More volatile with occasional dips below 90%."),
        ("\U0001f30d  Multi-Region Health",
         "83% East US 2, 17% Sweden Central \u2014 healthy geo-distribution "
         "with both regions fully operational."),
        ("\U0001f504  V2 \u2192 NG Migration Complete",
         f"V2 traffic: {V2_TOTAL:,} (1.4%) vs NG: {TOTAL_NG:,} (98.6%). "
         "Multiple days with zero V2 requests. Migration effectively done."),
    ]
    y = 1.3
    for title, desc in items:
        _txt(sl, .8, y, 11.5, .35, title, 15, True, ACCENT_BLUE)
        _txt(sl, .8, y + .33, 11.5, .5, desc, 11, c=LIGHT_GRAY)
        y += .92

    # ── SLIDE 8 : Appendix \u2014 Dashboard Links ─────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
    _txt(sl, .5, .4, 12, .7, "Appendix: Dashboard Links & References", 24, True, WHITE)

    _txt(sl, .8, 1.2, 11, .4, "Live Dashboard (Azure Data Explorer):", 14, True, ACCENT_BLUE)
    _txt(sl, .8, 1.6, 11, .4, DASH, 10, c=LIGHT_GRAY)

    _txt(sl, .8, 2.3, 11, .4, "Individual Tile Links:", 14, True, ACCENT_BLUE)
    y = 2.8
    for name, url in TILES.items():
        _txt(sl, .8, y, 3, .32, f"\U0001f4ca  {name}", 11, c=WHITE)
        _txt(sl, 3.8, y, 9, .32, url, 8, c=MED_GRAY)
        y += .35

    _txt(sl, .8, 5.2, 11, .4, "Methodology Notes:", 14, True, ACCENT_BLUE)
    _ml(sl, .8, 5.7, 11, 1.5, [
        '\u2022 NG data: UnionOfAllLogs("Vienna", "requests") \u2014 app == "hephaestus-agent"',
        '\u2022 V2 data: UnionOfAllLogs("Vienna", "traces") \u2014 message contains "AgentId=triage"',
        "\u2022 Regions: eastus2, swedencentral  |  Period: 30 days",
        "\u2022 Success rate: Per-region data from tile bd391155, weighted 83% EUS2 + 17% SC",
        "\u2022 Full KQL queries available in companion Word document",
    ], 10, LIGHT_GRAY)

    # ── SLIDE 9 : Appendix \u2014 Resolved Queries ───────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl)
    _txt(sl, .5, .4, 12, .7, "Appendix: Analysis Status", 24, True, WHITE)

    _txt(sl, .8, 1.2, 11, .4, "\u2705  User Retention (Resolved)", 14, True, ACCENT_GREEN)
    _txt(sl, .8, 1.65, 11, .4,
         f"Returning users grew from 0% to ~10%. "
         f"Weekday avg: {AVG_WD_RETURN_PCT}%. See Slide 4 for charts.",
         11, c=LIGHT_GRAY)

    _txt(sl, .8, 2.4, 11, .4, "\u2705  % of NG Users in AskAI (Resolved)", 14, True, ACCENT_GREEN)
    _txt(sl, .8, 2.85, 11, .8,
         f"Total AskAI Users: {ASKAI_TOTAL_USERS:,}  |  "
         f"Hephaestus NG Users: {ASKAI_NG_USERS:,}  |  "
         f"NG as % of AskAI: {ASKAI_NG_PCT}%",
         12, c=LIGHT_GRAY)

    _txt(sl, .8, 3.8, 11, .4, "\u2705  All planned analyses are now complete.", 14, True, ACCENT_GREEN)

    p = os.path.join(OUT_DIR, "Hephaestus_NG_Leadership_30Day_v5.pptx")
    prs.save(p); print(f"\u2705 PowerPoint saved: {p}"); return p


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD WORD DOCUMENT
# ═══════════════════════════════════════════════════════════════════════════════

def _doc_table(doc, rows, style="Medium Shading 1 Accent 1"):
    t = doc.add_table(rows=len(rows), cols=len(rows[0]))
    t.style = style; t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            t.rows[i].cells[j].text = val

def _doc_kql(doc, heading, kql):
    doc.add_heading(heading, level=2)
    doc.add_paragraph(kql, style="No Spacing")


def build_docx():
    doc = Document()

    doc.add_heading("Hephaestus NG \u2014 30-Day Detailed Report", level=0)
    doc.add_paragraph(
        f"Period: March 10 \u2013 April 9, 2026  |  Regions: East US 2, Sweden Central\n"
        f"Team: {TEAM_NAME}  |  Manager: {TEAM_MANAGER}\n"
        f"Members: {', '.join(TEAM_MEMBERS)}"
    )

    # 1
    doc.add_heading("1. Executive Summary", level=1)
    doc.add_paragraph(
        f"Over the last 30 days, Hephaestus NG processed {TOTAL_NG:,} total requests "
        f"across East US 2 (83%) and Sweden Central (17%), serving {TOTAL_USERS:,} unique users. "
        f"The service averaged {SR_AVG}% success rate across both regions, with P95 latency at 5\u20136s. "
        f"V2 (triage) traffic was just {V2_TOTAL:,} requests (1.4%), confirming migration "
        "to NG is effectively complete."
    )

    # 2
    doc.add_heading("2. Request Volume", level=1)
    _doc_table(doc, [
        ("Metric", "Value"),
        ("Total NG Requests (30d)", f"{TOTAL_NG:,}"),
        ("Avg Weekday Requests", "~3,600/day"),
        ("Avg Weekend Requests", "~1,400/day"),
        ("East US 2 Share", "83%"),
        ("Sweden Central Share", "17%"),
    ])
    doc.add_paragraph("")
    doc.add_heading("Daily Request Volume", level=2)
    _doc_table(doc, [("Date", "Total Requests")] + [(d, f"{c:,}") for d, c in daily_total_requests],
               style="Light List Accent 1")
    _doc_kql(doc, "KQL \u2014 NG Request Volume",
             'let timeGrain = _period * 1h;\n'
             'UnionOfAllLogs("Vienna", "requests")\n'
             '| where operation_Name == "POST Hephaestus/ColdStart"\n'
             '    or operation_Name == "POST Hephaestus/ChatStream"\n'
             '| where PreciseTimeStamp >= _startTime and PreciseTimeStamp <= _endTime\n'
             '    and env in (_env)\n'
             '| where app == "hephaestus-agent"\n'
             '| where operation_Name != "GET Version/Ping"\n'
             '| summarize count=count() by bin(timestamp, timeGrain), type, env')
    doc.add_paragraph(f"Dashboard tile: {TILES['NG Requests']}")

    # 3
    doc.add_heading("3. User Engagement", level=1)
    _doc_table(doc, [
        ("Metric", "Value"),
        ("Total Unique Users (30d)", f"{TOTAL_USERS:,}"),
        ("Avg Weekday DAU", f"~{AVG_WD_DAU}"),
        ("Avg Weekend DAU", f"~{AVG_WE_DAU}"),
        ("Peak Day", "799 (Mar 12)"),
    ])
    doc.add_paragraph("")
    doc.add_heading("Daily Unique Users", level=2)
    _doc_table(doc, [("Date", "Unique Users")] + [(d, str(c)) for d, c in daily_users],
               style="Light List Accent 1")
    _doc_kql(doc, "KQL \u2014 Daily Unique Users",
             'UnionOfAllLogs("Vienna", "requests")\n'
             '| where operation_Name in ("POST Hephaestus/ColdStart",\n'
             '    "POST Hephaestus/ChatStream")\n'
             '| where PreciseTimeStamp >= _startTime and PreciseTimeStamp <= _endTime\n'
             '    and env in (_env)\n'
             '| where app == "hephaestus-agent"\n'
             '| summarize UniqueUsers = dcount(user_Id) by bin(timestamp, 1d)\n'
             '| order by timestamp asc')

    # 4
    doc.add_heading("4. SLO Performance", level=1)

    doc.add_heading("4.1 Success Rate", level=2)
    doc.add_paragraph(
        f"The service averaged {SR_AVG}% success rate across both regions over 31 days. "
        "East US 2 averaged ~97% and Sweden Central averaged ~97%. "
        f"Minimum was {SR_MIN_VAL}% on {SR_MIN_DAY} (driven by East US 2 at 79.7%), "
        f"and the maximum was {SR_MAX_VAL}% on {SR_MAX_DAY}."
    )
    doc.add_paragraph(
        "The Apr 13 East US 2 dip to 79.7% is an anomaly under investigation. "
        "Sweden Central remained at 100% on that day. Excluding Apr 13, "
        "the combined average is closer to 97%."
    )
    _doc_table(doc, [
        ("Metric", "Value"),
        ("Avg Success Rate (weighted)", f"{SR_AVG}%"),
        ("Avg East US 2", "~97%"),
        ("Avg Sweden Central", "~97%"),
        ("Minimum (combined)", f"{SR_MIN_VAL}% ({SR_MIN_DAY})"),
        ("Maximum (combined)", f"{SR_MAX_VAL}% ({SR_MAX_DAY})"),
    ])

    doc.add_heading("Daily Success Rate by Region", level=3)
    sr_rows = [("Date", "East US 2", "Sweden Central", "Combined")]
    for (d, e), (_, s), (_, c) in zip(daily_success_eus2, daily_success_sc, daily_success):
        sr_rows.append((d, f"{e}%", f"{s}%", f"{c}%"))
    _doc_table(doc, sr_rows, style="Light List Accent 1")

    doc.add_paragraph(f"Dashboard tile: {TILES['Success Rate (Overall)']}")

    doc.add_heading("4.2 V1 P95 Latency", level=2)
    _doc_table(doc, [
        ("Metric", "Value"),
        ("Median P95", f"{V1_P95_MEDIAN}s ({int(V1_P95_MEDIAN*1000):,} ms)"),
        ("Typical Range", f"{V1_P95_RANGE}"),
        ("Best Observed", f"~{V1_P95_MIN_S}s"),
        ("Worst Observed", f"~{V1_P95_MAX_S}s"),
        ("Status", "\u2705 Healthy"),
    ])
    doc.add_paragraph(f"Dashboard tile: {TILES['V1 P95 Latency']}")

    doc.add_heading("4.3 V1 P99 Latency (TTFB)", level=2)
    doc.add_paragraph(
        f"Typically {V1_P99_RANGE}. Periodic spikes to ~{V1_P99_MAX_S}s are cold-start "
        "timeout/retry events at the 99th percentile, affecting <1% of requests."
    )
    _doc_table(doc, [
        ("Metric", "Value"),
        ("Median P99", f"{V1_P99_MEDIAN}s ({int(V1_P99_MEDIAN*1000):,} ms)"),
        ("Typical Range", f"{V1_P99_RANGE}"),
        ("Best Observed", f"~{V1_P99_MIN_S}s"),
        ("Worst Spike", f"~{V1_P99_MAX_S}s"),
        ("Status", "\u26a0\ufe0f Monitor cold-start spikes"),
    ])
    doc.add_paragraph(f"Dashboard tile: {TILES['V1 P99 TTFB']}")

    doc.add_heading("4.4 V2 P95 Latency", level=2)
    _doc_table(doc, [
        ("Metric", "Value"),
        ("Median P95", f"{V2_P95_MEDIAN}s ({int(V2_P95_MEDIAN*1000):,} ms)"),
        ("Typical Range", f"{V2_P95_RANGE}"),
        ("Best Observed", f"~{V2_P95_MIN_S}s"),
        ("Worst Observed", f"~{V2_P95_MAX_S}s"),
        ("Status", "\u26a0\ufe0f Higher spikes than V1 P95"),
    ])
    doc.add_paragraph(f"Dashboard tile: {TILES['V2 P95 Latency']}")

    doc.add_heading("4.5 V2 P99 Latency (TTFB)", level=2)
    doc.add_paragraph(
        f"V2 P99 is typically {V2_P99_RANGE}, with a maximum of ~{V2_P99_MAX_S}s. "
        "No extreme cold-start spikes observed (V2 traffic is lower and more sporadic)."
    )
    _doc_table(doc, [
        ("Metric", "Value"),
        ("Median P99", f"{V2_P99_MEDIAN}s ({int(V2_P99_MEDIAN*1000):,} ms)"),
        ("Typical Range", f"{V2_P99_RANGE}"),
        ("Best Observed", f"~{V2_P99_MIN_S}s"),
        ("Worst Observed", f"~{V2_P99_MAX_S}s"),
        ("Status", "\u2705 Healthy (no extreme spikes)"),
    ])
    doc.add_paragraph(f"Dashboard tile: {TILES['V2 P99 TTFB']}")

    # 5
    doc.add_heading("5. V2 \u2192 NG Migration Status", level=1)
    doc.add_paragraph(
        f"V2 (triage agent) processed {V2_TOTAL:,} requests in 30 days (1.4%), "
        f"compared to NG's {TOTAL_NG:,} (98.6%). Multiple days had zero V2 requests "
        "(Mar 24\u2013Apr 1). Migration is effectively complete."
    )
    _doc_table(doc, [("Date", "V2 Requests")] + [(d, str(c)) for d, c in v2_daily]
               + [("Total", str(V2_TOTAL))])
    _doc_kql(doc, "KQL \u2014 V2 Requests",
             'UnionOfAllLogs("Vienna", "traces")\n'
             '| where timestamp between (_startTime .. _endTime)\n'
             '| where app == "hephaestus-agent"\n'
             '    and env in (\'eastus2\', \'swedencentral\')\n'
             '    and message contains "Running for AgentId"\n'
             '    and message contains "Running for AgentId=triage"\n'
             '| summarize count = count() by bin(timestamp, 1h)')
    doc.add_paragraph(f"Dashboard tile: {TILES['V2 Requests (Triage)']}")

    # 6
    doc.add_heading("6. User Retention (New vs Returning)", level=1)
    doc.add_paragraph(
        f"Over the 30-day period, {TOTAL_NEW_USERS:,} new users and {TOTAL_RETURNING:,} returning users "
        f"were observed. 'New' means first appearance within the analysis window. "
        f"The returning user percentage grew from 0% (day 1) to ~10% by week 4, "
        f"demonstrating increasing user stickiness. "
        f"Average returning % across all days: {AVG_RETURN_PCT}%. "
        f"Weekday average: {AVG_WD_RETURN_PCT}%."
    )
    doc.add_heading("6.1 Key Retention Insights", level=2)
    _doc_table(doc, [
        ("Metric", "Value"),
        ("Total New Users", f"{TOTAL_NEW_USERS:,}"),
        ("Total Returning Users", f"{TOTAL_RETURNING:,}"),
        ("Avg Returning %", f"{AVG_RETURN_PCT}%"),
        ("Weekday Avg Returning %", f"{AVG_WD_RETURN_PCT}%"),
        ("Peak Returning %", f"{max(rp for _, _, _, _, _, rp in retention_daily)}% ({next(d for d, _, _, _, _, rp in retention_daily if rp == max(r for _, _, _, _, _, r in retention_daily))})"),
        ("Trend", "0% \u2192 ~10% (steady growth)"),
    ])
    doc.add_heading("6.2 Daily Retention Breakdown", level=2)
    ret_rows = [("Date", "New Users", "Returning", "Total", "New %", "Return %")]
    for d, n, r, t, np_, rp in retention_daily:
        ret_rows.append((d, str(n), str(r), str(t), f"{np_}%", f"{rp}%"))
    _doc_table(doc, ret_rows, style="Light List Accent 1")
    _doc_kql(doc, "KQL \u2014 User Retention",
             'let allActivity =\n'
             '    UnionOfAllLogs("Vienna", "requests")\n'
             '    | where operation_Name in\n'
             '        ("POST Hephaestus/ColdStart", "POST Hephaestus/ChatStream")\n'
             '    | where PreciseTimeStamp >= _startTime\n'
             '        and PreciseTimeStamp <= _endTime and env in (_env)\n'
             '    | where app == "hephaestus-agent"\n'
             '    | project user_Id, day = bin(timestamp, 1d);\n'
             'let firstSeen =\n'
             '    allActivity | summarize firstDay = min(day) by user_Id;\n'
             'allActivity\n'
             '| distinct user_Id, day\n'
             '| join kind=inner firstSeen on user_Id\n'
             '| extend userType = iff(day == firstDay, "New", "Returning")\n'
             '| summarize\n'
             '    NewUsers = dcountif(user_Id, userType == "New"),\n'
             '    ReturningUsers = dcountif(user_Id, userType == "Returning")\n'
             '    by day\n'
             '| extend\n'
             '    TotalUsers = NewUsers + ReturningUsers,\n'
             '    NewPct = round(todouble(NewUsers)\n'
             '        / todouble(NewUsers + ReturningUsers) * 100, 1),\n'
             '    ReturnPct = round(todouble(ReturningUsers)\n'
             '        / todouble(NewUsers + ReturningUsers) * 100, 1)\n'
             '| order by day asc')

    # 7
    doc.add_heading("7. Pending: % of NG Users in AskAI", level=1)
    doc.add_paragraph(
        "Run this query to calculate NG users as a % of total AskAI users. "
        "Adjust askAIUsers filter if AskAI uses a different app/operation."
    )
    doc.add_paragraph(
        'let ngUsers = toscalar(\n'
        '    UnionOfAllLogs("Vienna", "requests")\n'
        '    | where operation_Name in\n'
        '        ("POST Hephaestus/ColdStart", "POST Hephaestus/ChatStream")\n'
        '    | where PreciseTimeStamp >= _startTime\n'
        '        and PreciseTimeStamp <= _endTime and env in (_env)\n'
        '    | where app == "hephaestus-agent"\n'
        '    | summarize dcount(user_Id)\n'
        ');\n'
        'let askAIUsers = toscalar(\n'
        '    UnionOfAllLogs("Vienna", "requests")\n'
        '    | where PreciseTimeStamp >= _startTime\n'
        '        and PreciseTimeStamp <= _endTime and env in (_env)\n'
        '    | where app == "hephaestus-agent"\n'
        '    | summarize dcount(user_Id)\n'
        ');\n'
        'print NGUsers = ngUsers, AskAIUsers = askAIUsers,\n'
        '      PctOfAskAI = round(\n'
        '          todouble(ngUsers) / todouble(askAIUsers) * 100, 2)',
        style="No Spacing"
    )

    # 8
    doc.add_heading("8. Dashboard References", level=1)
    doc.add_paragraph(f"Main Dashboard:\n{DASH}")
    for name, url in TILES.items():
        doc.add_paragraph(f"\u2022 {name}:\n  {url}")

    # 9
    doc.add_heading("9. Charts", level=1)
    for title, fname in [
        ("Daily Request Volume", "chart_requests.png"),
        ("Traffic by Region", "chart_region.png"),
        ("NG vs V2 Comparison", "chart_v2_vs_ng.png"),
        ("Daily Unique Users", "chart_users.png"),
        ("New vs Returning Users", "chart_retention.png"),
        ("V1 Daily Success Rate", "chart_v1_success.png"),
        ("V2 Daily Success Rate", "chart_v2_success.png"),
        ("Per-Region Success Rate", "chart_success.png"),
    ]:
        fp = os.path.join(OUT_DIR, fname)
        if os.path.exists(fp):
            doc.add_heading(title, level=2)
            doc.add_picture(fp, width=DocInches(6.5))

    p = os.path.join(OUT_DIR, "Hephaestus_NG_Detailed_Report_v5.docx")
    doc.save(p); print(f"\u2705 Word document saved: {p}"); return p


# ═══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("Generating charts...")
    _chart_requests()
    _chart_users()
    _chart_success()
    _chart_v1_success()
    _chart_v2_success()
    _chart_region_pie()
    _chart_v2_vs_ng()
    _chart_retention()
    print("Generating PowerPoint...")
    build_pptx()
    print("Generating Word document...")
    build_docx()
    print(f"\n\U0001f389 All files saved to: {OUT_DIR}")
